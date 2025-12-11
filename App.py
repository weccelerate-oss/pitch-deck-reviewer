"""
מערכת סקירת מצגות AI - גרסה סופית
Pitch Deck Review System with Native PowerPoint Comments

תמיכה בהערות PowerPoint מקוריות (Modern Comments)
עם ממשק עברית מתוקן
"""

import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from docx import Document
import google.generativeai as genai
import json
import io
import re
import zipfile
from datetime import datetime
from lxml import etree
import uuid
import copy


# ============================================================
# הגדרות API
# ============================================================
GEMINI_API_KEY = "AIzaSyBJstgLpy_6W8OkQTD6t8HmfYTLL1sTLXE"
genai.configure(api_key=GEMINI_API_KEY)


# ============================================================
# קבועים XML להערות PowerPoint מודרניות
# ============================================================
# Modern Comments (Office 2019+)
MODERN_COMMENTS_NS = "http://schemas.microsoft.com/office/powerpoint/2018/8/main"
MODERN_COMMENTS_REL_TYPE = "http://schemas.microsoft.com/office/2018/10/relationships/comments"
MODERN_COMMENTS_CONTENT_TYPE = "application/vnd.ms-powerpoint.comments+xml"

# Legacy Comments (Office 2007-2016)
LEGACY_COMMENTS_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
LEGACY_COMMENTS_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
LEGACY_COMMENTS_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
LEGACY_AUTHORS_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"
LEGACY_AUTHORS_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors"

# Namespaces
NSMAP_LEGACY = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


# ============================================================
# הגדרות עמוד Streamlit
# ============================================================
st.set_page_config(
    page_title="סקירת מצגות AI",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)


# ============================================================
# CSS מתקדם - עיצוב מרשים ומקצועי
# ============================================================
st.markdown("""
<style>
/* ===== ייבוא פונטים ===== */
@import url('https://fonts.googleapis.com/css2?family=Heebo:wght@300;400;500;600;700;800;900&display=swap');

/* ===== משתני עיצוב ===== */
:root {
    --bg-dark: #0f0f23;
    --bg-card: #1a1a2e;
    --bg-card-hover: #232342;
    --accent-primary: #6366f1;
    --accent-secondary: #8b5cf6;
    --accent-tertiary: #a855f7;
    --accent-success: #10b981;
    --accent-warning: #f59e0b;
    --accent-danger: #ef4444;
    --accent-info: #3b82f6;
    --text-primary: #f8fafc;
    --text-secondary: #94a3b8;
    --text-muted: #64748b;
    --border-color: #334155;
    --glow-primary: rgba(99, 102, 241, 0.4);
    --glow-success: rgba(16, 185, 129, 0.4);
    --radius-sm: 8px;
    --radius-md: 12px;
    --radius-lg: 16px;
    --radius-xl: 24px;
    --shadow-sm: 0 2px 8px rgba(0, 0, 0, 0.3);
    --shadow-md: 0 4px 20px rgba(0, 0, 0, 0.4);
    --shadow-lg: 0 8px 40px rgba(0, 0, 0, 0.5);
    --shadow-glow: 0 0 30px var(--glow-primary);
}

/* ===== איפוס והגדרות גלובליות ===== */
html, body, .stApp {
    direction: rtl !important;
    text-align: right !important;
    font-family: 'Heebo', -apple-system, BlinkMacSystemFont, sans-serif !important;
    background: linear-gradient(135deg, var(--bg-dark) 0%, #16162e 50%, #1a1a3e 100%) !important;
    color: var(--text-primary) !important;
}

* {
    line-height: 1.7 !important;
}

/* ===== תוכן ראשי ===== */
.main .block-container {
    direction: rtl !important;
    text-align: right !important;
    padding: 2rem 3rem !important;
    max-width: 1400px !important;
    background: transparent !important;
}

/* ===== סרגל צד מעוצב ===== */
[data-testid="stSidebar"] {
    direction: rtl !important;
    background: linear-gradient(180deg, #0c0c1d 0%, #12122a 50%, #0a0a1a 100%) !important;
    border-left: 1px solid var(--border-color) !important;
}

[data-testid="stSidebar"]::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 3px;
    background: linear-gradient(90deg, var(--accent-primary), var(--accent-secondary), var(--accent-tertiary));
}

[data-testid="stSidebar"] > div:first-child {
    direction: rtl !important;
    padding: 2rem 1.5rem !important;
}

[data-testid="stSidebar"] * {
    direction: rtl !important;
    text-align: right !important;
}

[data-testid="stSidebar"] .stMarkdown p,
[data-testid="stSidebar"] .stMarkdown span,
[data-testid="stSidebar"] .stMarkdown h1,
[data-testid="stSidebar"] .stMarkdown h2,
[data-testid="stSidebar"] .stMarkdown h3 {
    color: var(--text-primary) !important;
}

[data-testid="stSidebar"] hr {
    border: none !important;
    height: 1px !important;
    background: linear-gradient(90deg, transparent, var(--border-color), transparent) !important;
    margin: 1.5rem 0 !important;
}

/* ===== מדדים בסרגל צד - תיקון קריאות ===== */
[data-testid="stSidebar"] [data-testid="stMetric"] {
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.15) 0%, rgba(139, 92, 246, 0.1) 100%) !important;
    border: 1px solid rgba(99, 102, 241, 0.3) !important;
    border-radius: var(--radius-md) !important;
    padding: 1rem !important;
    margin: 0.5rem 0 !important;
}

[data-testid="stSidebar"] [data-testid="stMetricLabel"] {
    color: var(--text-secondary) !important;
    font-size: 0.85rem !important;
    font-weight: 500 !important;
}

[data-testid="stSidebar"] [data-testid="stMetricValue"] {
    color: var(--text-primary) !important;
    font-size: 1.8rem !important;
    font-weight: 700 !important;
    text-shadow: 0 0 20px var(--glow-primary) !important;
}

/* ===== כותרות ===== */
h1, h2, h3, h4, h5, h6 {
    color: var(--text-primary) !important;
    font-weight: 700 !important;
    margin-bottom: 1rem !important;
    direction: rtl !important;
    text-align: right !important;
}

h1 { font-size: 2.5rem !important; }
h2 { font-size: 1.8rem !important; }
h3 { font-size: 1.4rem !important; }

/* ===== פסקאות וטקסט ===== */
p, span, div, label, li {
    direction: rtl !important;
    text-align: right !important;
    color: var(--text-secondary) !important;
}

.stMarkdown p {
    color: var(--text-secondary) !important;
}

/* ===== כרטיס כותרת ראשית ===== */
.main-header {
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.1) 0%, rgba(139, 92, 246, 0.05) 100%);
    border: 1px solid rgba(99, 102, 241, 0.2);
    border-radius: var(--radius-xl);
    padding: 3rem 2rem;
    margin-bottom: 2rem;
    text-align: center;
    position: relative;
    overflow: hidden;
}

.main-header::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 4px;
    background: linear-gradient(90deg, var(--accent-primary), var(--accent-secondary), var(--accent-tertiary), var(--accent-primary));
    background-size: 300% 100%;
    animation: gradient-flow 5s ease infinite;
}

@keyframes gradient-flow {
    0%, 100% { background-position: 0% 50%; }
    50% { background-position: 100% 50%; }
}

.main-title {
    font-size: 3rem !important;
    font-weight: 900 !important;
    background: linear-gradient(135deg, #fff 0%, #a5b4fc 50%, #c4b5fd 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin-bottom: 0.5rem !important;
    text-shadow: none;
    letter-spacing: -1px;
}

.sub-title {
    color: var(--text-secondary) !important;
    font-size: 1.2rem !important;
    font-weight: 400 !important;
    margin: 0 !important;
}

/* ===== שדות קלט ===== */
.stTextInput > div,
.stTextArea > div,
.stSelectbox > div {
    direction: rtl !important;
}

.stTextInput input,
.stTextArea textarea {
    direction: rtl !important;
    text-align: right !important;
    background: var(--bg-card) !important;
    border: 2px solid var(--border-color) !important;
    border-radius: var(--radius-md) !important;
    color: var(--text-primary) !important;
    padding: 0.875rem 1rem !important;
    font-family: 'Heebo', sans-serif !important;
    font-size: 1rem !important;
    transition: all 0.3s ease !important;
}

.stTextInput input:focus,
.stTextArea textarea:focus {
    border-color: var(--accent-primary) !important;
    box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.2), var(--shadow-glow) !important;
    outline: none !important;
}

.stTextInput input::placeholder,
.stTextArea textarea::placeholder {
    color: var(--text-muted) !important;
}

/* ===== תיבות בחירה ===== */
.stSelectbox [data-baseweb="select"] {
    direction: rtl !important;
}

.stSelectbox [data-baseweb="select"] > div {
    background: var(--bg-card) !important;
    border: 2px solid var(--border-color) !important;
    border-radius: var(--radius-md) !important;
    color: var(--text-primary) !important;
    direction: rtl !important;
    text-align: right !important;
}

.stSelectbox [data-baseweb="select"] > div:hover {
    border-color: var(--accent-primary) !important;
}

/* ===== העלאת קבצים ===== */
[data-testid="stFileUploader"] {
    direction: rtl !important;
}

[data-testid="stFileUploader"] section {
    direction: rtl !important;
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.05) 0%, rgba(139, 92, 246, 0.03) 100%) !important;
    border: 2px dashed var(--border-color) !important;
    border-radius: var(--radius-lg) !important;
    padding: 2.5rem !important;
    transition: all 0.3s ease !important;
}

[data-testid="stFileUploader"] section:hover {
    border-color: var(--accent-primary) !important;
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.1) 0%, rgba(139, 92, 246, 0.05) 100%) !important;
    box-shadow: var(--shadow-glow) !important;
}

[data-testid="stFileUploader"] section > div {
    text-align: center !important;
}

[data-testid="stFileUploader"] small {
    color: var(--text-muted) !important;
}

/* ===== כפתורים ===== */
.stButton > button {
    direction: rtl !important;
    font-family: 'Heebo', sans-serif !important;
    font-weight: 600 !important;
    font-size: 1rem !important;
    border-radius: var(--radius-md) !important;
    padding: 0.875rem 2rem !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    position: relative !important;
    overflow: hidden !important;
}

.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, var(--accent-primary) 0%, var(--accent-secondary) 100%) !important;
    border: none !important;
    color: white !important;
    box-shadow: 0 4px 15px rgba(99, 102, 241, 0.4) !important;
}

.stButton > button[kind="primary"]:hover {
    transform: translateY(-3px) !important;
    box-shadow: 0 8px 25px rgba(99, 102, 241, 0.5), var(--shadow-glow) !important;
}

.stButton > button[kind="primary"]:active {
    transform: translateY(-1px) !important;
}

.stButton > button[kind="secondary"] {
    background: transparent !important;
    border: 2px solid var(--border-color) !important;
    color: var(--text-primary) !important;
}

.stButton > button[kind="secondary"]:hover {
    border-color: var(--accent-primary) !important;
    background: rgba(99, 102, 241, 0.1) !important;
}

.stButton > button:disabled {
    background: var(--bg-card) !important;
    border: 2px solid var(--border-color) !important;
    color: var(--text-muted) !important;
    box-shadow: none !important;
    cursor: not-allowed !important;
    opacity: 0.6 !important;
}

/* ===== כפתורי הורדה ===== */
[data-testid="stDownloadButton"] > button {
    background: linear-gradient(135deg, var(--accent-success) 0%, #059669 100%) !important;
    border: none !important;
    color: white !important;
    font-weight: 600 !important;
    box-shadow: 0 4px 15px rgba(16, 185, 129, 0.4) !important;
}

[data-testid="stDownloadButton"] > button:hover {
    transform: translateY(-3px) !important;
    box-shadow: 0 8px 25px rgba(16, 185, 129, 0.5), 0 0 30px var(--glow-success) !important;
}

/* ===== הודעות מערכת ===== */
[data-testid="stAlert"] {
    direction: rtl !important;
    border-radius: var(--radius-md) !important;
    padding: 1rem 1.25rem !important;
    margin: 1rem 0 !important;
    border: none !important;
}

[data-testid="stAlert"] > div {
    direction: rtl !important;
    text-align: right !important;
}

/* הצלחה */
.stSuccess, [data-testid="stAlert"][data-baseweb*="positive"] {
    background: linear-gradient(135deg, rgba(16, 185, 129, 0.15) 0%, rgba(16, 185, 129, 0.05) 100%) !important;
    border-right: 4px solid var(--accent-success) !important;
}

.stSuccess p, [data-testid="stAlert"][data-baseweb*="positive"] p {
    color: #6ee7b7 !important;
}

/* אזהרה */
.stWarning {
    background: linear-gradient(135deg, rgba(245, 158, 11, 0.15) 0%, rgba(245, 158, 11, 0.05) 100%) !important;
    border-right: 4px solid var(--accent-warning) !important;
}

.stWarning p {
    color: #fcd34d !important;
}

/* מידע */
.stInfo {
    background: linear-gradient(135deg, rgba(59, 130, 246, 0.15) 0%, rgba(59, 130, 246, 0.05) 100%) !important;
    border-right: 4px solid var(--accent-info) !important;
}

.stInfo p {
    color: #93c5fd !important;
}

/* שגיאה */
.stError {
    background: linear-gradient(135deg, rgba(239, 68, 68, 0.15) 0%, rgba(239, 68, 68, 0.05) 100%) !important;
    border-right: 4px solid var(--accent-danger) !important;
}

.stError p {
    color: #fca5a5 !important;
}

/* ===== טבלה ועורך נתונים ===== */
[data-testid="stDataEditor"],
[data-testid="stDataFrame"] {
    direction: rtl !important;
    border-radius: var(--radius-lg) !important;
    overflow: hidden !important;
    box-shadow: var(--shadow-lg) !important;
    margin: 1.5rem 0 !important;
    border: 1px solid var(--border-color) !important;
}

[data-testid="stDataEditor"] > div,
[data-testid="stDataFrame"] > div {
    direction: rtl !important;
    background: var(--bg-card) !important;
}

/* כותרות טבלה */
[data-testid="stDataEditor"] [role="columnheader"],
[data-testid="stDataFrame"] [role="columnheader"] {
    direction: rtl !important;
    text-align: right !important;
    font-weight: 700 !important;
    padding: 1rem !important;
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.2) 0%, rgba(139, 92, 246, 0.1) 100%) !important;
    color: var(--text-primary) !important;
    border-bottom: 2px solid var(--accent-primary) !important;
    font-size: 0.95rem !important;
}

/* תאי טבלה */
[data-testid="stDataEditor"] [role="gridcell"],
[data-testid="stDataFrame"] [role="gridcell"] {
    direction: rtl !important;
    text-align: right !important;
    padding: 0.875rem 1rem !important;
    color: var(--text-secondary) !important;
    border-bottom: 1px solid var(--border-color) !important;
    background: var(--bg-card) !important;
}

[data-testid="stDataEditor"] [role="gridcell"]:hover,
[data-testid="stDataFrame"] [role="gridcell"]:hover {
    background: var(--bg-card-hover) !important;
}

/* שורות לסירוגין */
[data-testid="stDataEditor"] [role="row"]:nth-child(even) [role="gridcell"],
[data-testid="stDataFrame"] [role="row"]:nth-child(even) [role="gridcell"] {
    background: rgba(99, 102, 241, 0.03) !important;
}

/* ===== מדדים באזור הראשי ===== */
.main [data-testid="stMetric"] {
    background: linear-gradient(135deg, var(--bg-card) 0%, var(--bg-card-hover) 100%) !important;
    border: 1px solid var(--border-color) !important;
    border-radius: var(--radius-md) !important;
    padding: 1.25rem !important;
    box-shadow: var(--shadow-md) !important;
    transition: all 0.3s ease !important;
}

.main [data-testid="stMetric"]:hover {
    border-color: var(--accent-primary) !important;
    box-shadow: var(--shadow-lg), var(--shadow-glow) !important;
    transform: translateY(-2px) !important;
}

.main [data-testid="stMetricLabel"] {
    color: var(--text-secondary) !important;
    font-size: 0.9rem !important;
    font-weight: 500 !important;
}

.main [data-testid="stMetricValue"] {
    color: var(--text-primary) !important;
    font-size: 2rem !important;
    font-weight: 700 !important;
}

/* ===== אקספנדר ===== */
.streamlit-expanderHeader {
    direction: rtl !important;
    text-align: right !important;
    font-family: 'Heebo', sans-serif !important;
    font-weight: 600 !important;
    font-size: 1rem !important;
    background: linear-gradient(135deg, var(--bg-card) 0%, var(--bg-card-hover) 100%) !important;
    border: 1px solid var(--border-color) !important;
    border-radius: var(--radius-md) !important;
    padding: 1rem 1.5rem !important;
    color: var(--text-primary) !important;
    transition: all 0.3s ease !important;
}

.streamlit-expanderHeader:hover {
    border-color: var(--accent-primary) !important;
    background: var(--bg-card-hover) !important;
}

.streamlit-expanderContent {
    direction: rtl !important;
    text-align: right !important;
    background: var(--bg-card) !important;
    border: 1px solid var(--border-color) !important;
    border-top: none !important;
    border-radius: 0 0 var(--radius-md) var(--radius-md) !important;
    padding: 1.5rem !important;
}

.streamlit-expanderContent p {
    color: var(--text-secondary) !important;
    margin-bottom: 0.75rem !important;
}

/* ===== קו הפרדה ===== */
hr {
    border: none !important;
    height: 1px !important;
    background: linear-gradient(90deg, transparent 0%, var(--border-color) 20%, var(--accent-primary) 50%, var(--border-color) 80%, transparent 100%) !important;
    margin: 2.5rem 0 !important;
}

/* ===== Spinner ===== */
.stSpinner > div {
    direction: rtl !important;
    text-align: right !important;
    color: var(--text-secondary) !important;
}

/* ===== Caption ===== */
.stCaption, [data-testid="stCaptionContainer"] {
    direction: rtl !important;
    text-align: right !important;
    color: var(--text-muted) !important;
    font-size: 0.85rem !important;
}

/* ===== עמודות ===== */
[data-testid="column"] {
    padding: 0.75rem !important;
}

/* ===== אנימציות ===== */
@keyframes pulse-glow {
    0%, 100% { box-shadow: 0 0 20px rgba(99, 102, 241, 0.3); }
    50% { box-shadow: 0 0 40px rgba(99, 102, 241, 0.5); }
}

@keyframes float {
    0%, 100% { transform: translateY(0); }
    50% { transform: translateY(-5px); }
}

/* ===== פס גלילה ===== */
::-webkit-scrollbar {
    width: 10px;
    height: 10px;
}

::-webkit-scrollbar-track {
    background: var(--bg-dark);
    border-radius: 5px;
}

::-webkit-scrollbar-thumb {
    background: linear-gradient(135deg, var(--accent-primary) 0%, var(--accent-secondary) 100%);
    border-radius: 5px;
    border: 2px solid var(--bg-dark);
}

::-webkit-scrollbar-thumb:hover {
    background: linear-gradient(135deg, var(--accent-secondary) 0%, var(--accent-tertiary) 100%);
}

/* ===== התאמה למובייל ===== */
@media (max-width: 768px) {
    .main .block-container {
        padding: 1rem !important;
    }
    
    .main-title {
        font-size: 2rem !important;
    }
    
    .sub-title {
        font-size: 1rem !important;
    }
    
    .main-header {
        padding: 2rem 1rem !important;
    }
    
    [data-testid="stDataEditor"] {
        font-size: 0.85rem !important;
    }
    
    .stButton > button {
        width: 100% !important;
        padding: 1rem !important;
    }
    
    [data-testid="column"] {
        width: 100% !important;
        margin-bottom: 0.75rem !important;
    }
}

/* ===== תיקוני RTL נוספים ===== */
[data-testid="stDataEditor"] input,
[data-testid="stDataEditor"] select {
    direction: rtl !important;
    text-align: right !important;
}

/* Fix for select dropdowns */
[data-baseweb="popover"] {
    direction: rtl !important;
}

[data-baseweb="menu"] {
    direction: rtl !important;
}

[data-baseweb="menu"] li {
    direction: rtl !important;
    text-align: right !important;
}
</style>
""", unsafe_allow_html=True)


# ============================================================
# פונקציות חילוץ טקסט
# ============================================================

def extract_text_from_shape(shape) -> str:
    """חילוץ טקסט רקורסיבי מצורות."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    text_parts = []
    
    if hasattr(shape, "text_frame"):
        try:
            text = shape.text_frame.text.strip()
            if text:
                text_parts.append(text)
        except:
            pass
    elif hasattr(shape, "text"):
        try:
            text = shape.text.strip()
            if text:
                text_parts.append(text)
        except:
            pass
    
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            for row in shape.table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    text_parts.append(" | ".join(row_texts))
        except:
            pass
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                child_text = extract_text_from_shape(child)
                if child_text:
                    text_parts.append(child_text)
        except:
            pass
    
    if hasattr(shape, "has_chart") and shape.has_chart:
        try:
            if shape.chart.has_title:
                title = shape.chart.chart_title.text_frame.text.strip()
                if title:
                    text_parts.append(f"[תרשים: {title}]")
        except:
            pass
    
    return "\n".join(text_parts)


def extract_text_from_pptx(file_bytes: bytes) -> list[dict]:
    """חילוץ טקסט מכל השקפים."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(io.BytesIO(file_bytes))
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        has_visuals = False
        
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                texts.append(text)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_visuals = True
            if hasattr(shape, "has_chart") and shape.has_chart:
                has_visuals = True
        
        slide_text = "\n".join(texts).strip()
        
        if not slide_text:
            slide_text = "[שקף עם תמונות/גרפיקה]" if has_visuals else "[שקף ריק]"
        
        slides_data.append({"slide_number": slide_num, "text": slide_text})
    
    return slides_data


def extract_text_from_docx(file_bytes: bytes) -> str:
    """חילוץ טקסט מ-Word."""
    doc = Document(io.BytesIO(file_bytes))
    return "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_text_from_txt(file_bytes: bytes) -> str:
    """חילוץ טקסט מקובץ טקסט."""
    for enc in ["utf-8", "utf-8-sig", "windows-1255", "iso-8859-8", "latin-1"]:
        try:
            return file_bytes.decode(enc)
        except:
            continue
    return file_bytes.decode("utf-8", errors="replace")


# ============================================================
# ניתוח AI
# ============================================================

# System Prompt מקצועי לבדיקת מצגות
SYSTEM_PROMPT = """
תפקיד:
אתה משמש כבודק מצגות מקצועי. מטרתך היא לבחון מצגות סקירה של מיזמים ולספק הערות מדויקות, ממוקדות וברורות, תוך היצמדות מוחלטת לשיחת הפתיחה של המיזם (Context).

כללי יסוד:
1. לפני תחילת הבדיקה נדרשת שיחת פתיחה (Context).
2. אם קיים חוסר התאמה בין המצגת לשיחת הפתיחה (קהל יעד, תחום, סוג מוצר) — יש לציין זאת במפורש.
3. הקפד על: בהירות, רצף קריא, תמצות ודיוק.

הנחיות ניסוח להערות:
* כל הערה נפתחת כך: "שקף X – ..."
* ההנחיות חייבות להיות ישירות, לא כלליות.
* אין להוסיף ניסוחים ארוכים או חזרתיים.
* אין להעיר על שקף המתודולוגיה (התעלם ממנו).
* סדר הבדיקה: לפי הסדר במצגת, ללא מספור פנימי בתוך ההערה.

הנחיות תוכן ספציפיות:
* הגדרת הבעיה והפתרון: חייבים להופיע בשקפים נפרדים. ניסוח בפסקה קצרה וזורמת (לא בולטים).
    - בעיה: בלי ניסוחים מוחלטים ("לא קיים פתרון"), זהירות עם פתרונות חלקיים.
    - פתרון: תיאור המענה ללא השוואה למתחרים כרגע.
* נתוני שוק: מותר להציג רק נתונים מ-2024–2025. נתונים ישנים יותר -> חובה להעיר ולדרוש עדכון.
* מתחרים: מתחרים ישירים תחילה. חובה לכלול מתחרים ישראליים אם יש. לציין תאריך לנתונים (למשל "נכון ל-2019").
* סקירת שווקים: רק שווקים רלוונטיים. אם יש עומס - להעיר מה להסיר.
* מסקנות: חובה להציג חסרונות מתחרים + כיוון בידול ברור.

פלט נדרש:
החזר אך ורק JSON תקני המכיל רשימה של אובייקטים, כאשר כל אובייקט מכיל:
- "slide_number": מספר השקופית (מספר שלם).
- "original_text": תקציר קצר של תוכן השקף (עד 100 תווים).
- "ai_comment": ההערה המקצועית שלך לפי הכללים הנ"ל. אם השקף תקין, כתוב "תקין".
- "status": אחד מהבאים:
    * "לביצוע" - יש בעיה שדורשת תיקון
    * "אהבתי" - השקף מצוין, ראוי לציון חיובי
    * "נפתר" - השקף תקין, אין הערות
"""


def analyze_slides(slides_data: list[dict], context_text: str, model_name: str = "gemini-2.0-flash") -> list[dict]:
    """ניתוח שקפים עם Gemini AI באמצעות System Prompt מקצועי."""
    
    total = len(slides_data)
    
    # הכנת תוכן השקפים
    slides_content = "\n\n".join([
        f"=== שקף {s['slide_number']}/{total} ===\n{s['text']}"
        for s in slides_data
    ])
    
    # בניית הפרומפט המלא
    user_prompt = f"""
שיחת פתיחה (Context):
---
{context_text}
---

מצגת לבדיקה ({total} שקפים):
---
{slides_content}
---

בצע בדיקה מקצועית לכל {total} השקפים והחזר JSON בלבד.
"""

    try:
        model = genai.GenerativeModel(
            model_name,
            system_instruction=SYSTEM_PROMPT
        )
        
        response = model.generate_content(
            user_prompt,
            generation_config=genai.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.2,
                max_output_tokens=16384
            )
        )
        
        response_text = response.text.strip()
        
        # ניסיון לפענח JSON
        try:
            result = json.loads(response_text)
        except json.JSONDecodeError as je:
            # ניסיון לתקן JSON לא שלם
            st.warning(f"⚠️ תשובת AI לא תקינה, מנסה לתקן...")
            
            # הסרת backticks אם יש
            if response_text.startswith("```"):
                response_text = re.sub(r'^```json?\s*', '', response_text)
                response_text = re.sub(r'\s*```$', '', response_text)
            
            try:
                result = json.loads(response_text)
            except:
                st.error(f"❌ לא ניתן לפענח תשובת AI")
                st.code(response_text[:500], language="json")
                result = []
                
    except Exception as e:
        error_msg = str(e)
        
        if "429" in error_msg or "quota" in error_msg.lower() or "resource" in error_msg.lower():
            st.error("❌ חריגה ממכסת API - נסה מאוחר יותר או החלף מפתח")
        elif "403" in error_msg and "leaked" in error_msg.lower():
            st.error("❌ מפתח ה-API דווח כחשוף - צור מפתח חדש ב-Google AI Studio")
        elif "404" in error_msg:
            st.error(f"❌ מודל '{model_name}' לא נמצא - נסה מודל אחר")
        elif "API_KEY" in error_msg.upper() or "invalid" in error_msg.lower():
            st.error("❌ מפתח API לא תקין - בדוק את המפתח בקוד")
        else:
            st.error(f"❌ שגיאת API: {error_msg}")
        
        result = []
    
    # מילוי שקפים חסרים
    returned = {r.get("slide_number") for r in result}
    for slide in slides_data:
        if slide["slide_number"] not in returned:
            result.append({
                "slide_number": slide["slide_number"],
                "original_text": slide["text"][:100],
                "ai_comment": "⚠️ לא נותח - יש לסקור ידנית",
                "status": "לביצוע"
            })
    
    return sorted(result, key=lambda x: x.get("slide_number", 0))


# ============================================================
# הערות PowerPoint מקוריות - מניפולציית ZIP/XML
# ============================================================

def escape_xml(text: str) -> str:
    """Escape special XML characters."""
    return (text
            .replace('&', '&amp;')
            .replace('<', '&lt;')
            .replace('>', '&gt;')
            .replace('"', '&quot;')
            .replace("'", '&apos;'))


def create_comment_authors_xml() -> str:
    """יצירת XML של מחברי הערות."""
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cmAuthor id="1" name="AI Reviewer" initials="AI" lastIdx="1000" clrIdx="0"/>
</p:cmAuthorLst>'''


def create_slide_comments_xml(comments: list[dict]) -> str:
    """יצירת XML של הערות לשקף."""
    comments_xml = []
    for c in comments:
        dt = datetime.now().strftime("%Y-%m-%dT%H:%M:%S.000")
        text = escape_xml(c['text'])
        comments_xml.append(f'''
    <p:cm authorId="1" dt="{dt}" idx="{c['idx']}">
        <p:pos x="{c['x']}" y="{c['y']}"/>
        <p:text>{text}</p:text>
    </p:cm>''')
    
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">{"".join(comments_xml)}
</p:cmLst>'''


def add_comments_via_zip(pptx_bytes: bytes, analyzed_data: list[dict]) -> tuple[bytes, str]:
    """
    הוספת הערות PowerPoint מקוריות באמצעות מניפולציית ZIP ישירה.
    """
    debug_info = []
    
    try:
        # סינון הערות לפי סטטוס
        comments_by_slide = {}
        for item in analyzed_data:
            slide_num = item.get("slide_number")
            status = item.get("status", "")
            comment = item.get("ai_comment", "").strip()
            
            # דילוג על הערות "לא נותח"
            if "לא נותח" in comment:
                continue
            
            if status in ["נפתר", "למחוק"] or not comment:
                continue
            
            prefix = "[לביצוע]" if status == "לביצוע" else "[אהבתי]"
            full_comment = f"{prefix} {comment}"
            
            if slide_num not in comments_by_slide:
                comments_by_slide[slide_num] = []
            comments_by_slide[slide_num].append(full_comment)
        
        debug_info.append(f"שקפים עם הערות: {list(comments_by_slide.keys())}")
        
        if not comments_by_slide:
            return pptx_bytes, "⚠️ אין הערות להוספה (כל ההערות בסטטוס נפתר/למחוק)"
        
        # פתיחת ה-PPTX כ-ZIP
        input_zip = zipfile.ZipFile(io.BytesIO(pptx_bytes), 'r')
        output_buffer = io.BytesIO()
        output_zip = zipfile.ZipFile(output_buffer, 'w', zipfile.ZIP_DEFLATED)
        
        # רשימת כל הקבצים ב-ZIP
        all_files = input_zip.namelist()
        debug_info.append(f"קבצים ב-PPTX: {len(all_files)}")
        
        # קריאת [Content_Types].xml
        content_types = input_zip.read('[Content_Types].xml').decode('utf-8')
        
        # קריאת presentation.xml.rels
        pres_rels_path = 'ppt/_rels/presentation.xml.rels'
        pres_rels = input_zip.read(pres_rels_path).decode('utf-8')
        
        # בדיקה אם commentAuthors קיים
        has_authors = 'commentAuthors.xml' in content_types
        
        # הוספת commentAuthors אם לא קיים
        if not has_authors:
            insert_pos = content_types.rfind('</Types>')
            new_type = '<Override PartName="/ppt/commentAuthors.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"/>\n'
            content_types = content_types[:insert_pos] + new_type + content_types[insert_pos:]
            
            rid_matches = re.findall(r'Id="rId(\d+)"', pres_rels)
            max_rid = max([int(r) for r in rid_matches]) if rid_matches else 0
            new_rid = f"rId{max_rid + 1}"
            
            insert_pos = pres_rels.rfind('</Relationships>')
            new_rel = f'<Relationship Id="{new_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors" Target="commentAuthors.xml"/>\n'
            pres_rels = pres_rels[:insert_pos] + new_rel + pres_rels[insert_pos:]
            debug_info.append("נוצר commentAuthors.xml")
        
        # מעקב אחר קבצי הערות שנוצרו
        comments_files_added = []
        slide_rels_to_create = {}  # rels חדשים ליצירה
        slide_rels_to_update = {}  # rels קיימים לעדכון
        
        comment_idx = 1
        for slide_num, comments in comments_by_slide.items():
            # יצירת רשימת הערות
            comment_list = []
            for i, text in enumerate(comments):
                comment_list.append({
                    'idx': comment_idx,
                    'text': text,
                    'x': 7000000,
                    'y': 500000 + (i * 1200000)
                })
                comment_idx += 1
            
            # יצירת XML להערות
            comments_xml = create_slide_comments_xml(comment_list)
            comments_filename = f'ppt/comments/comment{slide_num}.xml'
            comments_files_added.append((comments_filename, comments_xml))
            
            # הוספה ל-Content_Types
            insert_pos = content_types.rfind('</Types>')
            new_type = f'<Override PartName="/{comments_filename}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.comments+xml"/>\n'
            content_types = content_types[:insert_pos] + new_type + content_types[insert_pos:]
            
            # בדיקה אם קיים rels לשקף
            slide_rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
            
            if slide_rels_path in all_files:
                slide_rels_to_update[slide_rels_path] = slide_num
                debug_info.append(f"שקף {slide_num}: יעודכן rels קיים")
            else:
                # יצירת rels חדש
                slide_rels_to_create[slide_rels_path] = slide_num
                debug_info.append(f"שקף {slide_num}: ייווצר rels חדש")
        
        # העתקת כל הקבצים עם עדכונים
        for item in all_files:
            if item == '[Content_Types].xml':
                output_zip.writestr(item, content_types.encode('utf-8'))
            elif item == pres_rels_path:
                output_zip.writestr(item, pres_rels.encode('utf-8'))
            elif item in slide_rels_to_update:
                # עדכון slide rels קיים
                slide_num = slide_rels_to_update[item]
                slide_rels = input_zip.read(item).decode('utf-8')
                
                rid_matches = re.findall(r'Id="rId(\d+)"', slide_rels)
                max_rid = max([int(r) for r in rid_matches]) if rid_matches else 0
                new_rid = f"rId{max_rid + 1}"
                
                insert_pos = slide_rels.rfind('</Relationships>')
                new_rel = f'<Relationship Id="{new_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" Target="../comments/comment{slide_num}.xml"/>\n'
                slide_rels = slide_rels[:insert_pos] + new_rel + slide_rels[insert_pos:]
                
                output_zip.writestr(item, slide_rels.encode('utf-8'))
            else:
                output_zip.writestr(item, input_zip.read(item))
        
        # יצירת rels חדשים לשקפים שאין להם
        for rels_path, slide_num in slide_rels_to_create.items():
            new_rels = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" Target="../comments/comment{slide_num}.xml"/>
</Relationships>'''
            output_zip.writestr(rels_path, new_rels.encode('utf-8'))
        
        # הוספת commentAuthors.xml אם לא קיים
        if not has_authors:
            output_zip.writestr('ppt/commentAuthors.xml', create_comment_authors_xml().encode('utf-8'))
        
        # הוספת קבצי הערות
        for filename, xml_content in comments_files_added:
            output_zip.writestr(filename, xml_content.encode('utf-8'))
        
        input_zip.close()
        output_zip.close()
        
        output_buffer.seek(0)
        
        # Debug output
        debug_str = " | ".join(debug_info)
        return output_buffer.getvalue(), f"✅ נוספו {len(comments_files_added)} הערות ({debug_str})"
        
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        return pptx_bytes, f"❌ שגיאה: {str(e)}"


def add_comments_to_speaker_notes(pptx_bytes: bytes, analyzed_data: list[dict]) -> tuple[bytes, int]:
    """Fallback: הוספת הערות ל-Speaker Notes."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    added_count = 0
    
    for item in analyzed_data:
        slide_num = item.get("slide_number", 0)
        status = item.get("status", "")
        comment = item.get("ai_comment", "").strip()
        
        # דילוג על הערות לא רלוונטיות
        if "לא נותח" in comment:
            continue
        if status in ["נפתר", "למחוק"] or not comment:
            continue
        
        if slide_num < 1 or slide_num > len(prs.slides):
            continue
        
        slide = prs.slides[slide_num - 1]
        
        indicator = "🔴 לביצוע" if status == "לביצוע" else "💚 אהבתי"
        formatted = f"\n\n{'='*40}\n{indicator} | AI Reviewer:\n{comment}\n{'='*40}"
        
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        existing = tf.text or ""
        
        tf.clear()
        tf.paragraphs[0].text = existing + formatted
        added_count += 1
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue(), added_count


def add_comments_to_pptx(pptx_bytes: bytes, analyzed_data: list[dict]) -> tuple[bytes, str]:
    """הוספת הערות - ניסיון מקוריות עם fallback ל-Speaker Notes."""
    
    # ניסיון ראשון: הערות מקוריות
    result_bytes, message = add_comments_via_zip(pptx_bytes, analyzed_data)
    
    if message.startswith("✅"):
        return result_bytes, message
    
    # Fallback: Speaker Notes
    try:
        fallback_bytes, count = add_comments_to_speaker_notes(pptx_bytes, analyzed_data)
        if count > 0:
            return fallback_bytes, f"⚠️ נוספו {count} הערות ל-Speaker Notes (הערות מקוריות: {message})"
        else:
            return pptx_bytes, f"⚠️ לא נוספו הערות ({message})"
    except Exception as e:
        return pptx_bytes, f"❌ שגיאה: {e}"


def create_excel_report(analyzed_data: list[dict]) -> bytes:
    """יצירת דוח Excel."""
    df = pd.DataFrame(analyzed_data)
    df = df.rename(columns={
        "slide_number": "מספר שקף",
        "original_text": "טקסט מקורי",
        "ai_comment": "הערת AI",
        "status": "סטטוס"
    })
    
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="ניתוח")
        ws = writer.sheets["ניתוח"]
        for i, col in enumerate(df.columns):
            ws.column_dimensions[chr(65 + i)].width = min(max(df[col].astype(str).str.len().max(), len(col)) + 2, 50)
    
    output.seek(0)
    return output.getvalue()


# ============================================================
# אפליקציה ראשית
# ============================================================

def main():
    # כותרת ראשית מעוצבת
    st.markdown('''
    <div class="main-header">
        <h1 class="main-title">🎯 מערכת סקירת מצגות AI</h1>
        <p class="sub-title">נתח את המצגת שלך באמצעות בינה מלאכותית וקבל הערות מקצועיות ישירות לתוך PowerPoint</p>
    </div>
    ''', unsafe_allow_html=True)
    
    # מדריך
    with st.expander("❓ איך להשתמש?", expanded=False):
        st.markdown("""
        ### 🚀 מדריך מהיר
        
        **1️⃣ העלאת קבצים**
        
        העלה מצגת PowerPoint וקובץ הקשר (TXT/DOCX)
        
        **2️⃣ ניתוח**
        
        לחץ על "נתח מצגת" והמתן לסיום
        
        **3️⃣ עריכה**
        
        סקור את ההערות בטבלה ועדכן סטטוסים
        
        **4️⃣ הורדה**
        
        הורד מצגת עם הערות בפאנל Review
        
        ---
        
        ### 📊 סטטוסים
        
        | סמל | סטטוס | יתווסף? |
        |:---:|:------|:-------:|
        | ⏳ | לביצוע | ✅ |
        | ❤️ | אהבתי | ✅ |
        | ✅ | נפתר | ❌ |
        | 🗑️ | למחוק | ❌ |
        
        ---
        
        💡 ההערות יופיעו ב-**Review > Comments** בפאוורפוינט
        """)
    
    st.markdown("---")
    
    # סרגל צד
    with st.sidebar:
        st.markdown("## ⚙️ הגדרות")
        st.markdown("")
        
        # API
        st.markdown("### 🔐 חיבור API")
        api_ok = GEMINI_API_KEY and GEMINI_API_KEY != "YOUR_API_KEY_HERE"
        if api_ok:
            st.markdown("✅ **מחובר** לשרת Gemini")
        else:
            st.markdown("❌ **לא מחובר** - עדכן מפתח API")
        
        st.markdown("---")
        
        # מודל
        st.markdown("### 🤖 מודל AI")
        model = st.selectbox(
            "בחירה",
            ["gemini-2.0-flash", "gemini-1.5-pro-latest", "gemini-1.5-flash-latest"],
            label_visibility="collapsed"
        )
        model_desc = {
            "gemini-2.0-flash": "⚡ **מהיר** - מומלץ לרוב המשימות", 
            "gemini-1.5-pro-latest": "🎯 **מדויק** - לניתוח מעמיק", 
            "gemini-1.5-flash-latest": "🚀 **קל** - לניתוח מהיר"
        }
        st.markdown(model_desc.get(model, ""))
        
        st.markdown("---")
        
        # סטטיסטיקות מעוצבות
        st.markdown("### 📊 סטטיסטיקות")
        
        if "slides_data" in st.session_state and st.session_state["slides_data"]:
            slides_count = len(st.session_state["slides_data"])
            st.markdown(f"""
            <div style="
                background: linear-gradient(135deg, rgba(99, 102, 241, 0.2) 0%, rgba(139, 92, 246, 0.1) 100%);
                border: 1px solid rgba(99, 102, 241, 0.3);
                border-radius: 12px;
                padding: 1rem;
                margin: 0.5rem 0;
                text-align: center;
            ">
                <div style="font-size: 2.5rem; font-weight: 700; color: #a5b4fc;">{slides_count}</div>
                <div style="font-size: 0.9rem; color: #94a3b8;">שקפים במצגת</div>
            </div>
            """, unsafe_allow_html=True)
        
        if "analysis_results" in st.session_state and st.session_state["analysis_results"]:
            df = pd.DataFrame(st.session_state["analysis_results"])
            counts = df["status"].value_counts()
            
            # סטטיסטיקות בכרטיסים
            st.markdown(f"""
            <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 0.5rem; margin-top: 1rem;">
                <div style="
                    background: linear-gradient(135deg, rgba(245, 158, 11, 0.2) 0%, rgba(245, 158, 11, 0.05) 100%);
                    border: 1px solid rgba(245, 158, 11, 0.3);
                    border-radius: 10px;
                    padding: 0.75rem;
                    text-align: center;
                ">
                    <div style="font-size: 1.8rem; font-weight: 700; color: #fcd34d;">⏳ {counts.get("לביצוע", 0)}</div>
                    <div style="font-size: 0.75rem; color: #94a3b8;">לביצוע</div>
                </div>
                <div style="
                    background: linear-gradient(135deg, rgba(16, 185, 129, 0.2) 0%, rgba(16, 185, 129, 0.05) 100%);
                    border: 1px solid rgba(16, 185, 129, 0.3);
                    border-radius: 10px;
                    padding: 0.75rem;
                    text-align: center;
                ">
                    <div style="font-size: 1.8rem; font-weight: 700; color: #6ee7b7;">✅ {counts.get("נפתר", 0)}</div>
                    <div style="font-size: 0.75rem; color: #94a3b8;">נפתר</div>
                </div>
                <div style="
                    background: linear-gradient(135deg, rgba(236, 72, 153, 0.2) 0%, rgba(236, 72, 153, 0.05) 100%);
                    border: 1px solid rgba(236, 72, 153, 0.3);
                    border-radius: 10px;
                    padding: 0.75rem;
                    text-align: center;
                ">
                    <div style="font-size: 1.8rem; font-weight: 700; color: #f9a8d4;">❤️ {counts.get("אהבתי", 0)}</div>
                    <div style="font-size: 0.75rem; color: #94a3b8;">אהבתי</div>
                </div>
                <div style="
                    background: linear-gradient(135deg, rgba(239, 68, 68, 0.2) 0%, rgba(239, 68, 68, 0.05) 100%);
                    border: 1px solid rgba(239, 68, 68, 0.3);
                    border-radius: 10px;
                    padding: 0.75rem;
                    text-align: center;
                ">
                    <div style="font-size: 1.8rem; font-weight: 700; color: #fca5a5;">🗑️ {counts.get("למחוק", 0)}</div>
                    <div style="font-size: 0.75rem; color: #94a3b8;">למחוק</div>
                </div>
            </div>
            """, unsafe_allow_html=True)
    
    # העלאת קבצים
    st.markdown("## 📂 העלאת קבצים")
    st.markdown("")
    
    c1, c2 = st.columns(2)
    
    with c1:
        st.markdown("#### 📑 מצגת")
        pptx_file = st.file_uploader("PPTX", type=["pptx"], key="pptx", label_visibility="collapsed")
        if pptx_file:
            st.success(f"✅ {pptx_file.name}")
    
    with c2:
        st.markdown("#### 💬 הקשר")
        context_file = st.file_uploader("TXT/DOCX", type=["txt", "docx"], key="ctx", label_visibility="collapsed")
        if context_file:
            st.success(f"✅ {context_file.name}")
    
    # עיבוד קבצים
    slides_data = context_text = None
    
    if pptx_file:
        pptx_file.seek(0)
        pptx_bytes = pptx_file.read()
        st.session_state["pptx_bytes"] = pptx_bytes
        slides_data = extract_text_from_pptx(pptx_bytes)
        st.session_state["slides_data"] = slides_data
    
    if context_file:
        ctx_bytes = context_file.read()
        context_text = extract_text_from_docx(ctx_bytes) if context_file.name.endswith(".docx") else extract_text_from_txt(ctx_bytes)
        st.session_state["context_text"] = context_text
    
    slides_data = slides_data or st.session_state.get("slides_data")
    context_text = context_text or st.session_state.get("context_text")
    
    st.markdown("---")
    
    # ניתוח
    st.markdown("## 🔬 ניתוח AI")
    st.markdown("")
    
    missing = []
    if not api_ok: missing.append("🔑 API")
    if not slides_data: missing.append("📑 מצגת")
    if not context_text: missing.append("💬 הקשר")
    
    can_analyze = not missing
    
    if missing:
        st.warning(f"⚠️ חסר: {' • '.join(missing)}")
    else:
        st.success("✅ מוכן לניתוח!")
    
    _, btn_col, _ = st.columns([1, 2, 1])
    with btn_col:
        if st.button("🔬 נתח מצגת", disabled=not can_analyze, type="primary", use_container_width=True):
            with st.spinner("⏳ מנתח... (עד דקה)"):
                try:
                    results = analyze_slides(slides_data, context_text, model)
                    st.session_state["analysis_results"] = results
                    
                    # בדיקה אם הניתוח הצליח
                    successful = sum(1 for r in results if "לא נותח" not in r.get("ai_comment", ""))
                    total = len(results)
                    
                    if successful == total:
                        st.success(f"🎉 הושלם! {total} שקפים נותחו בהצלחה")
                        st.balloons()
                    elif successful > 0:
                        st.warning(f"⚠️ נותחו {successful} מתוך {total} שקפים")
                    else:
                        st.error("❌ הניתוח נכשל - בדוק את מפתח ה-API")
                        
                except Exception as e:
                    st.error(f"❌ שגיאה: {e}")
    
    st.markdown("---")
    
    # תוצאות
    if st.session_state.get("analysis_results"):
        st.markdown("## 📋 תוצאות")
        st.markdown("")
        
        df = pd.DataFrame(st.session_state["analysis_results"]).sort_values("slide_number").reset_index(drop=True)
        
        # סידור עמודות - סטטוס מימין (ראשון ב-RTL)
        column_order = ["status", "slide_number", "ai_comment", "original_text"]
        df = df[[col for col in column_order if col in df.columns]]
        
        edited = st.data_editor(
            df,
            use_container_width=True,
            hide_index=True,
            column_config={
                "status": st.column_config.SelectboxColumn(
                    "📊 סטטוס", 
                    options=["לביצוע", "נפתר", "אהבתי", "למחוק"], 
                    width="small",
                    help="בחר סטטוס להערה"
                ),
                "slide_number": st.column_config.NumberColumn(
                    "🔢 שקף", 
                    disabled=True, 
                    width="small"
                ),
                "ai_comment": st.column_config.TextColumn(
                    "💬 הערת AI", 
                    width="large",
                    help="ניתן לערוך את ההערה"
                ),
                "original_text": st.column_config.TextColumn(
                    "📄 טקסט מקורי", 
                    disabled=True, 
                    width="medium"
                ),
            },
            column_order=column_order
        )
        
        st.session_state["analysis_results"] = edited.to_dict("records")
        
        st.markdown("---")
        
        # הורדה
        st.markdown("## ⬇️ הורדה")
        st.markdown("")
        
        # שמירת שינויים מהעורך
        st.session_state["analysis_results"] = edited.to_dict("records")
        
        counts = edited["status"].value_counts()
        active = counts.get("לביצוע", 0) + counts.get("אהבתי", 0)
        
        # הצגת סיכום
        st.info(f"📊 **{active}** הערות יתווספו למצגת (לביצוע: {counts.get('לביצוע', 0)}, אהבתי: {counts.get('אהבתי', 0)})")
        
        # Debug expander
        with st.expander("🔍 Debug - מידע טכני"):
            st.write("סטטוסים:")
            st.write(dict(counts))
            st.write("דוגמת נתונים:")
            if st.session_state.get("analysis_results"):
                st.write(st.session_state["analysis_results"][:2])
        
        c1, c2, c3 = st.columns(3)
        
        with c1:
            if st.session_state.get("pptx_bytes"):
                try:
                    result, msg = add_comments_to_pptx(
                        st.session_state["pptx_bytes"], 
                        st.session_state["analysis_results"]
                    )
                    st.caption(msg)
                    st.download_button(
                        "📊 מצגת + הערות", 
                        result, 
                        "מצגת_עם_הערות.pptx", 
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True, 
                        type="primary"
                    )
                except Exception as e:
                    st.error(f"❌ {e}")
                    import traceback
                    st.code(traceback.format_exc())
            else:
                st.warning("⚠️ לא נמצאה מצגת מקורית")
        
        with c2:
            try:
                excel = create_excel_report(st.session_state["analysis_results"])
                st.download_button("📑 Excel", excel, "צ'קליסט.xlsx",
                                  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                  use_container_width=True)
            except Exception as e:
                st.error(f"❌ {e}")
        
        with c3:
            st.download_button("🔧 JSON", 
                              json.dumps(st.session_state["analysis_results"], ensure_ascii=False, indent=2),
                              "ניתוח.json", "application/json", use_container_width=True)
    else:
        st.info("💡 העלה קבצים והפעל ניתוח כדי לראות תוצאות")


if __name__ == "__main__":
    main()
